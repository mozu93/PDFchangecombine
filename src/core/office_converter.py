"""
Officeファイル変換モジュール
要件定義書 F-102, F-103 Office文書のPDF変換実装
"""

import os
from pathlib import Path
from typing import Optional
import subprocess
import tempfile

# Office文書処理ライブラリ
try:
    from docx import Document  # python-docx
    from openpyxl import load_workbook  # Excel処理
    from pptx import Presentation  # python-pptx
    from reportlab.pdfgen import canvas
    from reportlab.lib.pagesizes import A4, letter
    from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer
    from reportlab.lib.styles import getSampleStyleSheet
    from reportlab.lib.units import inch
    from reportlab.pdfbase import pdfmetrics
    from reportlab.pdfbase.cidfonts import UnicodeCIDFont
except ImportError as e:
    # 開発時のみ警告、実際の運用時は requirements.txt でインストール必須
    print(f"Office変換ライブラリ未インストール: {e}")

from ..utils.logger import logger
from ..config import MAX_FILE_SIZE_MB


class OfficeConverter:
    """Office文書からPDF変換を行うクラス"""
    
    def __init__(self):
        self._setup_fonts()
        logger.info("Officeコンバーター初期化完了")
    
    def _setup_fonts(self):
        """PDF生成用フォント設定"""
        try:
            # 日本語フォント登録
            pdfmetrics.registerFont(UnicodeCIDFont('HeiseiMin-W3'))
            pdfmetrics.registerFont(UnicodeCIDFont('HeiseiKakuGo-W5'))
            self.font_name = 'HeiseiMin-W3'
            logger.info("日本語フォント設定完了")
        except Exception as e:
            self.font_name = 'Helvetica'  # フォールバック
            logger.warning(f"日本語フォント設定失敗、英語フォントを使用: {e}")
    
    def convert_to_pdf(self, input_path: str, output_path: str) -> bool:
        """
        OfficeファイルのPDF変換メイン処理
        
        Args:
            input_path: 入力ファイルパス
            output_path: 出力PDFパス
            
        Returns:
            bool: 変換成功時True
        """
        try:
            file_ext = Path(input_path).suffix.lower()
            
            if file_ext in ['.docx', '.doc']:
                return self._convert_word_to_pdf(input_path, output_path)
            elif file_ext in ['.xlsx', '.xls']:
                return self._convert_excel_to_pdf(input_path, output_path)
            elif file_ext in ['.pptx', '.ppt']:
                return self._convert_powerpoint_to_pdf(input_path, output_path)
            else:
                logger.error(f"未対応のOffice形式: {file_ext}")
                return False
        
        except Exception as e:
            logger.error(f"Office変換エラー: {input_path} - {str(e)}", exc_info=True)
            return False
    
    def _convert_word_to_pdf(self, input_path: str, output_path: str) -> bool:
        """Word文書のPDF変換"""
        try:
            # python-docxでWordファイル読み込み
            doc = Document(input_path)
            
            # PDF作成
            pdf_doc = SimpleDocTemplate(output_path, pagesize=A4)
            styles = getSampleStyleSheet()
            story = []
            
            # スタイル調整（日本語対応）
            title_style = styles['Title']
            title_style.fontName = self.font_name
            normal_style = styles['Normal']
            normal_style.fontName = self.font_name
            
            # パラグラフを順次追加
            for paragraph in doc.paragraphs:
                if paragraph.text.strip():
                    # タイトル判定（簡易）
                    if len(paragraph.text) < 50 and paragraph.runs and paragraph.runs[0].bold:
                        p = Paragraph(paragraph.text, title_style)
                    else:
                        p = Paragraph(paragraph.text, normal_style)
                    story.append(p)
                    story.append(Spacer(1, 0.1 * inch))
            
            pdf_doc.build(story)
            
            logger.info(f"Word変換完了: {Path(input_path).name}")
            return True
            
        except Exception as e:
            logger.error(f"Word変換エラー: {input_path} - {str(e)}")
            return False
    
    def _convert_excel_to_pdf(self, input_path: str, output_path: str) -> bool:
        """Excel文書のPDF変換"""
        try:
            # openpyxlでExcelファイル読み込み
            workbook = load_workbook(input_path, read_only=True)
            
            # PDF作成
            c = canvas.Canvas(output_path, pagesize=A4)
            width, height = A4
            
            margin = 50
            y_position = height - margin
            line_height = 20
            
            for sheet_name in workbook.sheetnames:
                sheet = workbook[sheet_name]
                
                # シート名を描画
                c.setFont(self.font_name, 16)
                c.drawString(margin, y_position, f"Sheet: {sheet_name}")
                y_position -= line_height * 2
                
                # データ行数制限（性能要件対応）
                max_rows = min(sheet.max_row, 100)  # 最大100行
                max_cols = min(sheet.max_column, 10)  # 最大10列
                
                c.setFont(self.font_name, 10)
                
                for row in range(1, max_rows + 1):
                    if y_position < margin:
                        c.showPage()
                        y_position = height - margin
                    
                    x_position = margin
                    row_text = ""
                    
                    for col in range(1, max_cols + 1):
                        cell = sheet.cell(row=row, column=col)
                        cell_value = str(cell.value) if cell.value is not None else ""
                        row_text += f"{cell_value}\\t"
                    
                    if row_text.strip():
                        c.drawString(x_position, y_position, row_text)
                    y_position -= line_height
                
                # シート間改ページ
                if sheet_name != workbook.sheetnames[-1]:
                    c.showPage()
                    y_position = height - margin
            
            c.save()
            
            logger.info(f"Excel変換完了: {Path(input_path).name}")
            return True
            
        except Exception as e:
            logger.error(f"Excel変換エラー: {input_path} - {str(e)}")
            return False
    
    def _convert_powerpoint_to_pdf(self, input_path: str, output_path: str) -> bool:
        """PowerPoint文書のPDF変換"""
        try:
            # python-pptxでPowerPointファイル読み込み
            presentation = Presentation(input_path)
            
            # PDF作成
            c = canvas.Canvas(output_path, pagesize=A4)
            width, height = A4
            margin = 50
            
            for slide_num, slide in enumerate(presentation.slides, 1):
                # スライド番号
                c.setFont(self.font_name, 14)
                c.drawString(margin, height - margin, f"Slide {slide_num}")
                
                y_position = height - margin - 40
                
                # スライド内のテキストシェイプを処理
                for shape in slide.shapes:
                    if hasattr(shape, 'text') and shape.text.strip():
                        # テキストを行分割
                        lines = shape.text.strip().split('\\n')
                        
                        for line in lines:
                            if y_position < margin:
                                break
                            
                            c.setFont(self.font_name, 12)
                            c.drawString(margin, y_position, line)
                            y_position -= 20
                
                # 次のスライドへ改ページ
                if slide_num < len(presentation.slides):
                    c.showPage()
            
            c.save()
            
            logger.info(f"PowerPoint変換完了: {Path(input_path).name}")
            return True
            
        except Exception as e:
            logger.error(f"PowerPoint変換エラー: {input_path} - {str(e)}")
            return False
    
    def _try_system_conversion(self, input_path: str, output_path: str) -> bool:
        """システム変換の試行（LibreOffice等）"""
        try:
            # LibreOfficeがインストールされている場合の変換
            cmd = [
                "libreoffice",
                "--headless",
                "--convert-to", "pdf",
                "--outdir", str(Path(output_path).parent),
                input_path
            ]
            
            result = subprocess.run(cmd, capture_output=True, text=True, timeout=30)
            
            if result.returncode == 0:
                logger.info(f"LibreOffice変換成功: {Path(input_path).name}")
                return True
            else:
                logger.warning(f"LibreOffice変換失敗: {result.stderr}")
                return False
        
        except (subprocess.TimeoutExpired, FileNotFoundError):
            logger.info("LibreOffice未検出、独自変換を使用")
            return False
        except Exception as e:
            logger.warning(f"システム変換エラー: {str(e)}")
            return False